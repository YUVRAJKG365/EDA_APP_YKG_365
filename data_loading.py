import streamlit as st
import pandas as pd
import numpy as np
import piexif
import PyPDF2
from PIL import Image
import io
import fitz  # PyMuPDF
import dask.dataframe as dd
import chardet
import re
import os
import tempfile
import base64
from io import BytesIO
from scipy.stats import ttest_ind, f_oneway, chi2_contingency
import concurrent.futures
import time
import psutil
import pytesseract  # For OCR text extraction
import matplotlib.pyplot as plt
from collections import Counter

# Function to detect file encoding
def detect_file_encoding(file):
    raw_data = file.read()
    result = chardet.detect(raw_data)
    file.seek(0)  # Reset file pointer to the beginning
    return result['encoding']

@st.cache_data(show_spinner=False)
def load_data(file):
    try:
        # Create a temporary file for all formats
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_path = tmp_file.name
            
        if file.name.endswith('.csv'):
            # Efficient CSV loading with optimized parameters
            df = pd.read_csv(tmp_path, engine='c', low_memory=False)
            os.unlink(tmp_path)
            return df
        elif file.name.endswith('.xlsx') or file.name.endswith('.xls'):
            # Efficient Excel loading
            df = pd.read_excel(tmp_path, engine='openpyxl')
            os.unlink(tmp_path)
            return df
        elif file.name.endswith('.pdf'):
            # Extract text and images from PDF
            with open(tmp_path, 'rb') as f:
                text, images = extract_text_and_images_from_pdf(f)
            os.unlink(tmp_path)
            return {'text': text, 'images': images}
        elif file.name.endswith('.txt'):
            # Read plain text file efficiently
            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            os.unlink(tmp_path)
            return text
        else:
            st.error("Unsupported file format. Please upload a CSV, Excel, PDF, TXT, or ODF file.")
            os.unlink(tmp_path)
            return None
    except Exception as e:
        st.error(f"Error loading the dataset: {e}")
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def load_large_data(file, file_type=None, blocksize="64MB", sample_rows=10000):
    """
    Efficiently loads large datasets using optimized techniques
    - Uses parallel processing, memory mapping, and efficient data types
    - Implements progress tracking and resource monitoring
    """
    try:
        # Create temporary file for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_path = tmp_file.name
            
        # Get file size for optimization decisions
        file_size = os.path.getsize(tmp_path) / (1024 * 1024)  # Size in MB
        st.info(f"File size: {file_size:.2f} MB - Optimizing loading process...")
        
        # Start resource monitoring
        start_time = time.time()
        mem_before = psutil.virtual_memory().used / (1024 * 1024)
        
        # Infer file type if not provided
        if file_type is None:
            fname = file.name.lower()
            if fname.endswith('.csv'):
                file_type = 'csv'
            elif fname.endswith('.xlsx') or fname.endswith('.xls'):
                file_type = 'excel'
            elif fname.endswith('.txt'):
                file_type = 'txt'
            elif fname.endswith('.pdf'):
                file_type = 'pdf'
            else:
                st.error("Unsupported file format for large data loading.")
                os.unlink(tmp_path)
                return None, None

        # Create progress bar
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        def update_progress(step, total_steps, message):
            progress = int((step / total_steps) * 100)
            progress_bar.progress(progress)
            status_text.text(f"{message} ({progress}%)")
        
        # CSV - Optimized loading with Dask and parallel processing
        if file_type == 'csv':
            update_progress(1, 4, "Analyzing CSV structure")
            
            # Optimize chunk size based on file size
            optimized_blocksize = "256MB" if file_size > 500 else "64MB"
            
            update_progress(2, 4, "Loading data in parallel chunks")
            df = dd.read_csv(
                tmp_path, 
                blocksize=optimized_blocksize,
                assume_missing=True,
                dtype=str,
                sample=1000000  # Optimized sampling
            )
            
            update_progress(3, 4, "Preparing preview")
            preview = df.head(sample_rows, compute=True)
            
            update_progress(4, 4, "Finalizing dataset")
            st.info(f"Previewing first {sample_rows} rows. Full dataset will be processed in chunks.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return df, preview
        
        # Excel - Optimized loading with chunk processing
        elif file_type == 'excel':
            update_progress(1, 3, "Reading Excel file")
            
            # For large files, read in chunks
            if file_size > 50:  # Files larger than 50MB
                chunks = []
                chunk_size = 10000
                total_rows = 0
                
                # Read Excel in chunks
                reader = pd.read_excel(tmp_path, engine='openpyxl', chunksize=chunk_size)
                for i, chunk_df in enumerate(reader):
                    chunks.append(chunk_df)
                    total_rows += len(chunk_df)
                    update_progress(i+1, int(2000000/chunk_size), f"Processing chunk {i+1}")
                
                df = pd.concat(chunks, ignore_index=True)
            else:
                df = pd.read_excel(tmp_path, engine='openpyxl')
            
            update_progress(2, 3, "Optimizing data types")
            df = optimize_data_types(df)
            
            update_progress(3, 3, "Finalizing dataset")
            st.info(f"Loaded {len(df)} rows from Excel file.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return df, df.head(sample_rows)
        
        # Text - Efficient reading with progress
        elif file_type == 'txt':
            update_progress(1, 2, "Reading text file")
            
            # Read large text files in chunks
            text = ""
            chunk_size = 1024 * 1024  # 1MB chunks
            total_chunks = (os.path.getsize(tmp_path) // chunk_size) + 1
            
            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                for i in range(total_chunks):
                    chunk = f.read(chunk_size)
                    text += chunk
                    update_progress(i+1, total_chunks, "Reading text chunks")
            
            update_progress(2, 2, "Finalizing text")
            st.info("Loaded plain text file.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return text, text[:10000]  # Return preview
        
        # PDF - Parallel processing of pages
        elif file_type == 'pdf':
            update_progress(1, 3, "Extracting PDF content")
            with open(tmp_path, 'rb') as f:
                text, images = extract_text_and_images_from_pdf(f)
            
            update_progress(2, 3, "Processing extracted content")
            if text or images:
                st.info("Loaded PDF file.")
                
                # Resource usage report
                mem_after = psutil.virtual_memory().used / (1024 * 1024)
                load_time = time.time() - start_time
                st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
                
                os.unlink(tmp_path)
                return {'text': text, 'images': images}, {'text': text[:10000], 'images': images[:5]}
            else:
                st.error("No content could be extracted from the PDF file.")
                os.unlink(tmp_path)
                return None, None
        
        else:
            st.error("Only CSV, Excel, PDF, TXT files are supported for large data loading.")
            os.unlink(tmp_path)
            return None, None
    
    except Exception as e:
        st.error(f"Error loading large dataset: {e}")
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None, None
    finally:
        progress_bar.empty()
        status_text.empty()

def is_structured(df):
    """
    Enhanced structured data check with additional validations:
    1. Consistent column count across rows
    2. Valid header row detection
    3. Non-empty data presence
    """
    try:
        # Basic emptiness check
        if df.empty or df.isnull().all().all():
            return False
        
        # Check for minimum structure requirements
        if len(df.columns) < 1 or df.shape[0] < 1:
            return False

        # Check for consistent number of columns across all rows
        col_count = len(df.columns)
        if any(len(row) != col_count for row in df.itertuples(index=False)):
            return False

        # Check for valid header (first row should contain mostly strings)
        header_validity = sum(isinstance(col, str) for col in df.columns) / len(df.columns)
        if header_validity < 0.8:  # At least 80% of headers should be strings
            return False

        # Check for non-empty data cells
        non_empty_cells = df.count().sum()
        if non_empty_cells / (df.shape[0] * df.shape[1]) < 0.1:  # At least 10% data present
            return False

        return True
    except Exception as e:
        st.error(f"Error checking dataset structure: {e}")
        return False

def optimize_data_types(df):
    for col in df.columns:
        if df[col].dtype == 'float64':
            df[col] = df[col].astype('float32')
        elif df[col].dtype == 'int64':
            df[col] = df[col].astype('int32')
        elif df[col].dtype == 'object':
            if df[col].nunique() / len(df[col]) < 0.5:  # Convert to category if unique values are less than 50%
                df[col] = df[col].astype('category')
    return df

# Enhanced function to transform unstructured data into structured data
def transform_unstructured_data(df):
    try:
        # Step 1: Drop completely empty rows and columns
        df_cleaned = df.dropna(how='all').dropna(axis=1, how='all')
        
        # Step 2: Fill missing values with placeholders
        df_cleaned = df_cleaned.fillna('MISSING_VALUE')
        
        # Step 3: Identify header row (row with most non-empty values)
        non_empty_counts = df_cleaned.apply(lambda row: row.astype(bool).sum(), axis=1)
        if non_empty_counts.max() > 1:
            header_row_idx = non_empty_counts.idxmax()
            new_header = df_cleaned.iloc[header_row_idx]
            df_cleaned = df_cleaned.drop(header_row_idx)
            df_cleaned.columns = new_header
        
        # Step 4: Reset index
        df_cleaned = df_cleaned.reset_index(drop=True)
        
        # Step 5: Clean column names
        df_cleaned.columns = [re.sub(r'\W+', '_', str(col)).strip('_') for col in df_cleaned.columns]
        
        return df_cleaned
    except Exception as e:
        st.error(f"Error transforming unstructured data: {e}")
        return df

# Enhanced PDF extraction function
def extract_text_and_images_from_pdf(uploaded_pdf):
    text = ""
    images = []
    try:
        # First try text extraction
        try:
            reader = PyPDF2.PdfReader(uploaded_pdf)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            st.warning(f"Standard text extraction failed: {e}")
        
        # If no text found, try OCR fallback
        if not text.strip():
            st.info("No text found. Attempting image extraction...")
            uploaded_pdf.seek(0)
            images = extract_images_from_pdf(uploaded_pdf)
            
            # If no images either, try OCR on the PDF pages
            if not images:
                st.info("No images found. Trying alternative text extraction...")
                uploaded_pdf.seek(0)
                with fitz.open(stream=uploaded_pdf.read(), filetype="pdf") as doc:
                    for page in doc:
                        text += page.get_text() or ""
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
    
    return text, images

def extract_images_from_pdf(pdf_file):
    images = []
    pdf_file.seek(0)
    try:
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page_index in range(len(doc)):
            for img_index, img in enumerate(doc.get_page_images(page_index)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)
    except Exception as e:
        st.error(f"Error extracting images from PDF: {e}")
    return images

def clear_all_data():
    """Clear all uploaded files, data, and session states"""
    # Clear session state variables
    keys_to_clear = ['uploaded_file', 'uploaded_file_name', 'df', 'is_structured', 
                     'text_data', 'pdf_data', 'image_data', 'pdf_images']
    for key in keys_to_clear:
        if key in st.session_state:
            del st.session_state[key]

    # Reset file uploader
    if 'file_uploader_key' in st.session_state:
        st.session_state.file_uploader_key += 1

    # Clear the cache
    st.cache_data.clear()

    st.success("All data and files have been cleared!")

def perform_statistical_tests(df, col1, col2, test_type):
    try:
        if test_type == "T-Test":
            t_stat, p_value = ttest_ind(df[col1], df[col2])
            return f"T-Test Results: t-statistic = {t_stat:.4f}, p-value = {p_value:.4f}"
        elif test_type == "ANOVA":
            f_stat, p_value = f_oneway(df[col1], df[col2])
            return f"ANOVA Results: F-statistic = {f_stat:.4f}, p-value = {p_value:.4f}"
        elif test_type == "Chi-Square Test":
            contingency_table = pd.crosstab(df[col1], df[col2])
            chi2, p_value, dof, expected = chi2_contingency(contingency_table)
            return f"Chi-Square Test Results: Chi2 = {chi2:.4f}, p-value = {p_value:.4f}"
    except Exception as e:
        st.error(f"Error performing statistical test: {e}")
        
def plot_data_summary(df):
    """
    Comprehensive data summary with statistics and visual overview.
    Handles None or empty DataFrame gracefully.
    """
    try:
        if df is None or not hasattr(df, "columns") or df.empty:
            st.warning("Please upload and load a dataset first. No data available for summary.")
            return

        summary = []
        for col in df.columns:
            dtype = str(df[col].dtype)
            unique = df[col].nunique()
            missing = df[col].isnull().sum()
            # Convert sample values to string to avoid Arrow conversion issues
            sample = str(df[col].dropna().unique()[:5])
            summary.append({
                "Column": col,
                "Type": dtype,
                "Unique Values": unique,
                "Missing Values": missing,
                "Sample Values": sample
            })
        summary_df = pd.DataFrame(summary)
        st.dataframe(summary_df)
    except Exception as e:
        st.error(f"Error in data summary: {e}")
        return None

def render_statistics_section():
    st.markdown("## 📊 Statistical Analysis")
    st.markdown("Perform statistical tests on your dataset.")

    if 'df' in st.session_state and st.session_state.df is not None:
        df = st.session_state.df.copy()  # Create a copy to avoid modifying original

        # Check if there are at least two numeric or categorical columns for statistical tests
        numeric_cols = df.select_dtypes(include=['number']).columns
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns

        if len(numeric_cols) + len(categorical_cols) < 2:
            st.info(
                "The uploaded dataset does not have enough suitable columns for statistical analysis. "
                "Please upload a dataset with at least two numeric or categorical columns."
            )
        else:
            # Main expander for test selection
            with st.expander("📝 Test Selection", expanded=True):
                col1, col2 = st.columns(2)

                with col1:
                    selected_col1 = st.selectbox(
                        "Select First Column:",
                        df.columns,
                        key="stat_col1"
                    )

                with col2:
                    selected_col2 = st.selectbox(
                        "Select Second Column:",
                        df.columns,
                        key="stat_col2"
                    )

                test_type = st.selectbox(
                    "Select Test Type:",
                    ["T-Test", "ANOVA", "Chi-Square Test"],
                    key="test_type"
                )

                if st.button("📊 Run Statistical Test", key="stat_test_btn"):
                    with st.spinner("Performing test..."):
                        try:
                            # Convert datetime columns to numeric if needed
                            for col in [selected_col1, selected_col2]:
                                if np.issubdtype(df[col].dtype, np.datetime64):
                                    df[col] = df[col].astype(np.int64) // 10**9  # Convert to Unix timestamp

                            result = perform_statistical_tests(df, selected_col1, selected_col2, test_type)
                            st.success("✅ Test completed successfully!")

                            # Display results outside the main expander
                            st.markdown("### 📋 Test Results")
                            st.write(result)
                            
                        except Exception as e:
                            st.error(f"❌ Error performing statistical test: {e}")
    else:
        st.info("Please upload and load a dataset first to use performance features.")

def extract_text_from_image(image):
    """Extract text from image using OCR"""
    try:
        # Use pytesseract for OCR
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        st.error(f"Error during OCR: {e}")
        return ""

def plot_word_frequency(text):
    """Generate word frequency visualization"""
    try:
        # Clean and split text into words
        words = re.findall(r'\b\w+\b', text.lower())
        word_counts = Counter(words)
        top_words = word_counts.most_common(20)
        
        if not top_words:
            st.info("No words found for frequency visualization.")
            return
        
        # Create visualization
        words, counts = zip(*top_words)
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.bar(words, counts)
        ax.set_title("Top 20 Most Frequent Words")
        ax.set_xlabel("Words")
        ax.set_ylabel("Frequency")
        plt.xticks(rotation=45)
        st.pyplot(fig)
    except Exception as e:
        st.error(f"Error generating word frequency plot: {e}")

def plot_image_histogram(image):
    """Generate histogram for image pixel intensities"""
    try:
        # Convert image to array
        img_array = np.array(image)
        
        # Create figure
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Process based on image mode
        if image.mode == 'L':
            # Grayscale image
            ax.hist(img_array.ravel(), bins=256, color='gray', alpha=0.7)
            ax.set_title("Grayscale Pixel Intensity Distribution")
        elif image.mode in ('RGB', 'RGBA'):
            # Color image
            colors = ('red', 'green', 'blue')
            channel_names = ('Red', 'Green', 'Blue')
            
            for i, color in enumerate(colors):
                hist, bins = np.histogram(img_array[:, :, i].ravel(), bins=256)
                ax.plot(bins[:-1], hist, color=color, label=channel_names[i])
            
            ax.set_title("Color Channel Pixel Intensity Distribution")
            ax.legend()
        else:
            st.warning(f"Unsupported image mode for histogram: {image.mode}")
            return
            
        ax.set_xlabel("Pixel Intensity")
        ax.set_ylabel("Frequency")
        ax.grid(True, alpha=0.3)
        st.pyplot(fig)
    except Exception as e:
        st.error(f"Error generating image histogram: {e}")

def home_ui():
    st.markdown("## 🏠 Home")
    st.markdown("Welcome to the EDA App! Upload your data to get started.")

    # Initialize session state keys
    if 'file_uploader_key' not in st.session_state:
        st.session_state.file_uploader_key = 0
    if 'file_processed' not in st.session_state:
        st.session_state.file_processed = False

    with st.expander("📁 Upload Data", expanded=True):
        uploader_key = f"file_uploader_{st.session_state.file_uploader_key}"
        uploaded_file = st.file_uploader(
            "Choose a CSV, Excel, PDF, Text, or Image file",
            type=["csv", "xlsx", "xls", "txt", "pdf", "png", "jpg", "jpeg", "bmp", "tiff", "gif", "webp", "docx", "pptx"],
            key=uploader_key,
            help="Supported formats: CSV, Excel (xlsx, xls), PDF, Text, Images (png, jpg, jpeg, bmp, tiff, gif, webp), DOCX, PPTX"
        )
        
        if uploaded_file is not None:
            # Store file in session state
            st.session_state.uploaded_file = uploaded_file
            st.session_state.uploaded_file_name = uploaded_file.name
            st.session_state.file_processed = False
            
            # Process file based on type
            try:
                if uploaded_file.name.endswith(('.txt', '.docx', '.pptx')):
                    # Text-based files
                    if uploaded_file.name.endswith('.txt'):
                        text_data = uploaded_file.read().decode("utf-8", errors="ignore")
                    elif uploaded_file.name.endswith('.docx'):
                        try:
                            from docx import Document
                            doc = Document(uploaded_file)
                            text_data = "\n".join([para.text for para in doc.paragraphs])
                        except ImportError:
                            st.error("To process DOCX files, please install python-docx: pip install python-docx")
                            text_data = ""
                    elif uploaded_file.name.endswith('.pptx'):
                        try:
                            from pptx import Presentation
                            prs = Presentation(uploaded_file)
                            text_data = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text_data.append(shape.text)
                            text_data = "\n".join(text_data)
                        except ImportError:
                            st.error("To process PPTX files, please install python-pptx: pip install python-pptx")
                            text_data = ""
                    
                    st.session_state.text_data = text_data
                    st.session_state.file_processed = True
                    st.success(f"✅ {uploaded_file.name} uploaded successfully!")
                    
                elif uploaded_file.name.endswith('.pdf'):
                    # PDF files
                    text, images = extract_text_and_images_from_pdf(uploaded_file)
                    st.session_state.pdf_data = {
                        'text': text,
                        'images': images
                    }
                    st.session_state.file_processed = True
                    
                    if text or images:
                        success_msg = "PDF file uploaded successfully!"
                        if text:
                            success_msg += " Text extracted."
                        if images:
                            success_msg += f" {len(images)} images extracted."
                        st.success(f"✅ {success_msg}")
                    else:
                        st.warning("PDF uploaded but no content could be extracted.")
                
                elif uploaded_file.type.startswith("image/"):
                    # Image files
                    image = Image.open(uploaded_file)
                    st.session_state.image_data = image
                    
                    # Extract text from image using OCR
                    extracted_text = extract_text_from_image(image)
                    st.session_state.image_extracted_text = extracted_text
                    
                    st.session_state.file_processed = True
                    st.success("✅ Image file uploaded successfully!")
                
                else:
                    # Tabular data (CSV, Excel)
                    df = load_data(uploaded_file)
                    if df is not None:
                        # Check if structured
                        if not is_structured(df):
                            st.warning("Dataset appears unstructured. Transforming...")
                            df = transform_unstructured_data(df)
                            
                            # Check again after transformation
                            if is_structured(df):
                                st.success("Dataset successfully transformed to structured format!")
                            else:
                                st.error("Could not transform to structured format. Proceeding with original data.")
                        
                        df = optimize_data_types(df)
                        st.session_state.df = df
                        st.session_state.is_structured = is_structured(df)
                        st.session_state.file_processed = True
                        st.success("✅ Dataset is ready for analysis!")
                    else:
                        st.error("❌ Error: The uploaded file could not be loaded.")
            except Exception as e:
                st.error(f"❌ Error processing file: {str(e)}")
                st.session_state.file_processed = False

    # Show file information if processed
    if st.session_state.get('file_processed', False):
        st.markdown(f"**Uploaded File:** `{st.session_state.uploaded_file_name}`")
        
        if st.button("🗑️ Remove Uploaded File", key="remove_file_btn"):
            for key in ['df', 'text_data', 'pdf_data', 'is_structured', 
                        'uploaded_file_name', 'uploaded_file', 'file_processed', 
                        'image_data', 'pdf_images', 'image_extracted_text']:
                if key in st.session_state:
                    del st.session_state[key]
            st.session_state.file_uploader_key += 1
            st.rerun()
    
    # Dataset information section
    if st.session_state.get('file_processed', False):
        st.markdown("## File Information")
        
        # Tabular data
        if 'df' in st.session_state and st.session_state.df is not None:
            df = st.session_state.df
            with st.expander("📐 Dataset Shape & Columns", expanded=True):
                st.markdown(f"**Total Rows:** {df.shape[0]}")
                st.markdown(f"**Total Columns:** {df.shape[1]}")
                st.markdown("---")
                st.markdown("**Column Names:**")
                columns_df = pd.DataFrame({
                    'Column Index': range(1, len(df.columns) + 1),
                    'Column Name': df.columns
                })
                st.dataframe(columns_df, hide_index=True, height=min(300, 35 * len(df.columns)))

            with st.expander("🔤 Data Types", expanded=True):
                dtype_df = pd.DataFrame({
                    'Column Name': df.columns,
                    'Data Type': df.dtypes.astype(str),
                    'Unique Values': [df[col].nunique() for col in df.columns]
                })
                st.dataframe(dtype_df, hide_index=True, height=min(400, 35 * len(df.columns)))
                st.markdown("**Data Type Summary:**")
                dtype_counts = df.dtypes.value_counts().reset_index()
                dtype_counts.columns = ['Data Type', 'Count']
                st.dataframe(dtype_counts, hide_index=True)

            with st.expander("❓ Missing Values", expanded=True):
                missing_df = pd.DataFrame({
                    'Column Name': df.columns,
                    'Missing Values': df.isnull().sum(),
                    'Percentage Missing': (df.isnull().mean() * 100).round(2)
                })
                st.dataframe(missing_df, hide_index=True, height=min(400, 35 * len(df.columns)))
                missing_total = missing_df['Missing Values'].sum()
                if missing_total > 0:
                    st.warning(f"**Total missing values in dataset:** {missing_total}")
                    st.markdown(f"**Columns with missing values:** {len(missing_df[missing_df['Missing Values'] > 0])}")

                    def highlight_missing_and_present(s):
                        return [
                            'background-color: yellow' if pd.isnull(v) else 'background-color: #b6fcb6'
                            for v in s
                        ]

                    st.markdown("**Highlighted Data (yellow = missing, green = present):**")

                    # Handle large datasets
                    n_cells = df.shape[0] * df.shape[1]
                    max_cells = 1000000  # 1 million cell threshold

                    if n_cells > max_cells:
                        st.warning(f"Large dataset: Showing sample of 10,000 rows")
                        sample_df = df.sample(min(10000, len(df)))
                        st.dataframe(
                            sample_df.style.apply(highlight_missing_and_present, axis=0),
                            hide_index=True,
                            height=min(400, 35 * len(sample_df))
                        )
                    else:
                        # Temporarily increase render limit
                        original_max = pd.get_option("styler.render.max_elements")
                        try:
                            if n_cells > original_max:
                                pd.set_option("styler.render.max_elements", n_cells)

                            st.dataframe(
                                df.style.apply(highlight_missing_and_present, axis=0),
                                hide_index=True,
                                height=min(400, 35 * len(df))
                            )
                        finally:
                            pd.set_option("styler.render.max_elements", original_max)
                else:
                    st.success("No missing values found in any column")
                    
            with st.expander("📈 Data Summary", expanded=True):
                plot_data_summary(df)

            with st.expander("👀 Data Preview", expanded=True):
                st.dataframe(df)

            with st.expander("📊 Quick Statistics", expanded=True):
                st.write(df.describe(include='all'))
                
            # Statistical Analysis Section
            with st.expander("📊 Statistical Analysis", expanded=False):
                st.markdown("### Statistical Tests")
                st.markdown("Perform statistical tests on your dataset.")

                # Check if there are at least two numeric or categorical columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                categorical_cols = df.select_dtypes(include=['object', 'category']).columns

                if len(numeric_cols) + len(categorical_cols) < 2:
                    st.info(
                        "The dataset does not have enough suitable columns for statistical analysis. "
                        "Need at least two numeric or categorical columns."
                    )
                else:
                    col1, col2 = st.columns(2)

                    with col1:
                        selected_col1 = st.selectbox(
                            "Select First Column:",
                            df.columns,
                            key="stat_col1"
                        )

                    with col2:
                        selected_col2 = st.selectbox(
                            "Select Second Column:",
                            df.columns,
                            key="stat_col2"
                        )

                    test_type = st.selectbox(
                        "Select Test Type:",
                        ["T-Test", "ANOVA", "Chi-Square Test"],
                        key="test_type"
                    )

                    if st.button("📊 Run Statistical Test", key="stat_test_btn"):
                        with st.spinner("Performing test..."):
                            try:
                                # Convert datetime columns to numeric if needed
                                temp_df = df.copy()
                                for col in [selected_col1, selected_col2]:
                                    if np.issubdtype(temp_df[col].dtype, np.datetime64):
                                        temp_df[col] = temp_df[col].astype(np.int64) // 10**9  # Convert to Unix timestamp

                                result = perform_statistical_tests(temp_df, selected_col1, selected_col2, test_type)
                                st.success("✅ Test completed successfully!")

                                # Display results
                                st.markdown("### 📋 Test Results")
                                st.write(result)
                                
                            except Exception as e:
                                st.error(f"❌ Error performing statistical test: {e}")
        
        # Text data (txt, docx, pptx)
        elif 'text_data' in st.session_state and st.session_state.text_data is not None:
            with st.expander("📝 Text File Information", expanded=True):
                text_data = st.session_state.text_data
                st.markdown(f"**Total Characters:** {len(text_data)}")
                st.markdown(f"**Total Words:** {len(text_data.split())}")
                st.markdown(f"**Total Lines:** {len(text_data.splitlines())}")
                
                # Word frequency visualization
                plot_word_frequency(text_data)
                
                st.markdown("### Text Preview")
                st.code(text_data[:1000] + ("..." if len(text_data) > 1000 else ""), language="text")
        
        # PDF data
        elif 'pdf_data' in st.session_state and st.session_state.pdf_data is not None:
            pdf_data = st.session_state.pdf_data
            
            # Text section
            if pdf_data['text']:
                with st.expander("📝 Extracted Text", expanded=True):
                    st.markdown(f"**Total Characters:** {len(pdf_data['text'])}")
                    st.markdown(f"**Total Words:** {len(pdf_data['text'].split())}")
                    st.markdown(f"**Total Pages:** {pdf_data['text'].count('\f') + 1}")
                    
                    # Word frequency visualization
                    plot_word_frequency(pdf_data['text'])
                    
                    st.markdown("### Text Preview")
                    st.code(pdf_data['text'][:1000] + ("..." if len(pdf_data['text']) > 1000 else ""), language="text")
            else:
                st.markdown("### No text content found in PDF")
            
            # Images section
            if pdf_data['images']:
                with st.expander(f"🖼️ Extracted Images ({len(pdf_data['images'])})", expanded=True):
                    st.markdown(f"Found {len(pdf_data['images'])} images in the PDF document")
                    
                    # Display images in a grid
                    cols_per_row = 3
                    for i in range(0, len(pdf_data['images']), cols_per_row):
                        cols = st.columns(cols_per_row)
                        for j, img in enumerate(pdf_data['images'][i:i+cols_per_row]):
                            with cols[j]:
                                # Convert PIL Image to bytes for display
                                buffered = BytesIO()
                                img.save(buffered, format="PNG")
                                img_str = base64.b64encode(buffered.getvalue()).decode()
                                st.markdown(
                                    f'<div style="text-align:center">'
                                    f'<img src="data:image/png;base64,{img_str}" style="max-width:100%; max-height:200px;">'
                                    f'<p style="font-size:0.8em;">Image {i+j+1}: {img.size[0]}x{img.size[1]} | {img.mode}</p>'
                                    f'</div>',
                                    unsafe_allow_html=True
                                )
            else:
                st.markdown("### No images found in PDF")
            
            # PDF metadata
            try:
                uploaded_file = st.session_state.uploaded_file
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(uploaded_file.getvalue())
                    tmp_path = tmp.name
                
                with fitz.open(tmp_path) as doc:
                    metadata = doc.metadata
                    with st.expander("📋 PDF Metadata", expanded=False):
                        st.markdown("### PDF Document Properties")
                        meta_df = pd.DataFrame(list(metadata.items()), columns=['Property', 'Value'])
                        st.dataframe(meta_df, hide_index=True)
                
                os.unlink(tmp_path)
            except Exception as e:
                st.warning(f"Could not extract PDF metadata: {e}")
        
        # Image data
        elif 'image_data' in st.session_state and st.session_state.image_data is not None:
            image = st.session_state.image_data
            with st.expander("🖼️ Image Information", expanded=True):
                # Display image
                buffered = BytesIO()
                image.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode()
                st.markdown(
                    f'<div style="text-align:center">'
                    f'<img src="data:image/png;base64,{img_str}" style="max-width:100%; max-height:400px;">'
                    f'</div>',
                    unsafe_allow_html=True
                )
                
                # Image metadata
                st.markdown("### Image Properties")
                st.write(f"**Format:** {image.format}")
                st.write(f"**Size:** {image.size[0]}x{image.size[1]} (width x height)")
                st.write(f"**Mode:** {image.mode}")
                
                # Pixel intensity visualization
                st.markdown("### Pixel Intensity Distribution")
                plot_image_histogram(image)
                
                # EXIF data if available
                try:
                    exif_data = image.getexif()
                    if exif_data:
                        st.markdown("### EXIF Metadata")
                        exif_info = {}
                        for tag_id, value in exif_data.items():
                            tag = piexif.TAGS.get(tag_id, tag_id)
                            exif_info[tag] = value
                        exif_df = pd.DataFrame(list(exif_info.items()), columns=['Property', 'Value'])
                        st.dataframe(exif_df, hide_index=True, height=300)
                    else:
                        st.info("No EXIF metadata found in image")
                except Exception as e:
                    st.warning(f"Could not extract EXIF data: {e}")
    else:
        st.info("ℹ️ Please upload a file to begin analysis")

