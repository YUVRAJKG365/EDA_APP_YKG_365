# extraction.py
import os
import io
import re
import pytesseract
from PIL import Image, ImageOps, ImageFilter, ImageStat
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
import zipfile
import tempfile
import streamlit as st
import numpy as np
import cv2
import matplotlib.pyplot as plt
from skimage import feature, exposure
from sklearn.decomposition import PCA
from sklearn.manifold import TSNE
from sklearn.cluster import KMeans
from deepface import DeepFace
from collections import Counter
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO

from utils.data_loader import display_file_info, handle_file_upload
from utils.session_state_manager import get_session_manager

# Configuration for OCR
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\yuvra\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"

class ImageAnalyzer:
    """Class to handle advanced image analysis and processing"""
    
    def __init__(self, image):
        self.image = image
        self.np_image = np.array(image)
        if len(self.np_image.shape) == 2:  # Grayscale
            self.np_image = np.stack((self.np_image,)*3, axis=-1)
        self.metadata = self._get_metadata()
    
    def _get_metadata(self):
        """Extract basic image metadata"""
        return {
            "format": self.image.format,
            "size": self.image.size,
            "mode": self.image.mode,
            "width": self.image.width,
            "height": self.image.height,
            "aspect_ratio": self.image.width / self.image.height,
            "file_size": len(self.image.tobytes())
        }
    
    def get_stats(self):
        """Calculate image statistics"""
        stats = ImageStat.Stat(self.image)
        if self.image.mode == 'RGB':
            return {
                "mean": stats.mean,
                "stddev": stats.stddev,
                "median": stats.median,
                "extrema": stats.extrema
            }
        else:
            return {
                "mean": stats.mean[0],
                "stddev": stats.stddev[0],
                "median": stats.median[0],
                "extrema": stats.extrema[0]
            }
    
    def get_histogram(self):
        """Generate RGB histogram"""
        if self.image.mode == 'L':
            plt.figure()
            plt.hist(np.array(self.image).ravel(), bins=256)
            plt.title('Grayscale Histogram')
            buf = BytesIO()
            plt.savefig(buf, format='png')
            plt.close()
            buf.seek(0)
            return Image.open(buf)
        else:
            plt.figure()
            for i, color in enumerate(['red', 'green', 'blue']):
                hist = cv2.calcHist([self.np_image], [i], None, [256], [0, 256])
                plt.plot(hist, color=color)
            plt.title('RGB Histogram')
            buf = BytesIO()
            plt.savefig(buf, format='png')
            plt.close()
            buf.seek(0)
            return Image.open(buf)
    
    def preprocess(self, operation, **kwargs):
        """Apply various preprocessing operations"""
        ops = {
            'grayscale': lambda img: img.convert('L'),
            'resize': lambda img: img.resize((kwargs.get('width', 256), kwargs.get('height', 256))),
            'rotate': lambda img: img.rotate(kwargs.get('angle', 90)),
            'flip': lambda img: img.transpose(Image.FLIP_LEFT_RIGHT),
            'mirror': lambda img: img.transpose(Image.FLIP_TOP_BOTTOM),
            'blur': lambda img: img.filter(ImageFilter.GaussianBlur(kwargs.get('radius', 2))),
            'sharpen': lambda img: img.filter(ImageFilter.SHARPEN),
            'edge_enhance': lambda img: img.filter(ImageFilter.EDGE_ENHANCE),
            'contrast': lambda img: ImageOps.autocontrast(img),
            'threshold': lambda img: img.point(lambda p: 255 if p > kwargs.get('threshold', 128) else 0),
            'canny': self._apply_canny
        }
        
        if operation in ops:
            return ops[operation](self.image)
        return self.image
    
    def _apply_canny(self, img):
        """Apply Canny edge detection"""
        gray = np.array(img.convert('L'))
        edges = feature.canny(gray, sigma=2)
        return Image.fromarray((edges * 255).astype('uint8'))
    
    def detect_faces(self):
        """Detect faces using DeepFace"""
        try:
            results = DeepFace.analyze(self.np_image, actions=['emotion', 'age', 'gender'], enforce_detection=False)
            return results
        except Exception as e:
            st.warning(f"Face detection failed: {str(e)}")
            return []
    
    def extract_text(self):
        """Extract text using OCR"""
        try:
            return pytesseract.image_to_string(self.image)
        except Exception as e:
            st.warning(f"OCR failed: {str(e)}")
            return ""

def extract_from_pdf(file_path, is_scanned=False):
    """
    Extract text and images from PDF files
    Supports both regular and scanned PDFs
    """
    text = ""
    images = []
    doc = fitz.open(file_path)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Extract text
        if is_scanned:
            # Convert PDF page to image and perform OCR
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text += pytesseract.image_to_string(img) + "\n"
        else:
            page_text = page.get_text()
            if page_text.strip() == "":
                # Fallback to OCR if no text found
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img) + "\n"
            else:
                text += page_text + "\n"
        
        # Extract embedded images
        for img_index, img_info in enumerate(page.get_images(full=True)):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            try:
                img = Image.open(io.BytesIO(image_bytes))
                images.append(img)
            except:
                st.warning(f"Couldn't process image {img_index} on page {page_num+1}")
    
    doc.close()
    return text, images

def extract_from_docx(file_path):
    """Extract text and images from DOCX files"""
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    images = []
    
    # Extract images by unzipping DOCX
    with zipfile.ZipFile(file_path) as z:
        for file in z.namelist():
            if re.match(r'word/media/.*\.(jpg|jpeg|png|gif|bmp|tiff)', file, re.I):
                img_data = z.read(file)
                try:
                    img = Image.open(io.BytesIO(img_data))
                    images.append(img)
                except:
                    st.warning(f"Couldn't process image: {file}")
    
    return text, images

def extract_from_pptx(file_path):
    """Extract text and images from PPTX files"""
    prs = Presentation(file_path)
    text = ""
    images = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
            
            if shape.shape_type == 13:  # 13 indicates picture shape
                img_data = shape.image.blob
                try:
                    img = Image.open(io.BytesIO(img_data))
                    images.append(img)
                except:
                    st.warning("Couldn't process an image in slide")
    
    return text, images

def extract_from_excel(file_path):
    """Extract data from Excel files (xlsx, xls)"""
    xl = pd.ExcelFile(file_path)
    text = ""
    
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        text += f"--- {sheet_name} ---\n"
        text += df.to_string() + "\n\n"
    
    return text, []  # Excel doesn't contain extractable images

def extract_from_csv(file_path):
    """Extract data from CSV files"""
    df = pd.read_csv(file_path)
    return df.to_string(), []  # CSVs don't contain images

def cluster_images(images):
    """Cluster similar images using feature extraction and dimensionality reduction"""
    if len(images) < 2:
        return []
    
    # Extract features (using simple color histograms for demo)
    features = []
    for img in images:
        np_img = np.array(img.convert('RGB'))
        hist = cv2.calcHist([np_img], [0, 1, 2], None, [8, 8, 8], [0, 256, 0, 256, 0, 256])
        hist = cv2.normalize(hist, hist).flatten()
        features.append(hist)
    
    # Reduce dimensions for visualization
    pca = PCA(n_components=2)
    reduced = pca.fit_transform(features)
    
    # Cluster
    kmeans = KMeans(n_clusters=min(3, len(images)), random_state=42)
    clusters = kmeans.fit_predict(features)
    
    return reduced, clusters

def universal_extraction(uploaded_filext):
    """
    Universal extraction function for all supported file types
    Returns extracted text and list of images
    """
    file_ext = os.path.splitext(uploaded_filext.name)[1].lower()
    text = ""
    images = []
    
    try:
        # Save uploaded file to temp location
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            tmp_file.write(uploaded_filext.getvalue())
            tmp_path = tmp_file.name
        
        # Process based on file type
        if file_ext == '.pdf':
            # First try regular extraction, fallback to OCR if needed
            text, images = extract_from_pdf(tmp_path, is_scanned=False)
            if not text.strip():
                text, images = extract_from_pdf(tmp_path, is_scanned=True)
        
        elif file_ext == '.docx':
            text, images = extract_from_docx(tmp_path)
        
        elif file_ext == '.pptx':
            text, images = extract_from_pptx(tmp_path)
        
        elif file_ext in ('.xlsx', '.xls'):
            text, images = extract_from_excel(tmp_path)
        
        elif file_ext == '.csv':
            text, images = extract_from_csv(tmp_path)
        
        else:
            st.error(f"Unsupported file type: {file_ext}")
            return "", []
        
        # Cleanup temporary file
        os.unlink(tmp_path)
        return text, images
    
    except Exception as e:
        st.error(f"Extraction failed: {str(e)}")
        return "", []

def show_image_analysis(images):
    tracker = st.session_state.tracker
    """Display advanced image analysis options for selected images"""
    if not images:
        st.warning("No images selected for analysis")
        return
    
    # Create tabs for analysis options
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "Preprocessing", "Statistics", 
        "OCR", "Visualizations", "Metadata"
    ])

    with tab1:
        tracker.log_tab("Preprocessing")
        st.subheader("Image Preprocessing")
        operation = st.selectbox(
            "Select operation",
            ["grayscale", "resize", "rotate", "flip", "mirror",
             "blur", "sharpen", "edge_enhance", "contrast", "threshold", "canny"],
            key="preprocess_operation"
        )
        

        params = {}
        if operation == "resize":
            params["width"] = st.slider("Width", 50, 1000, 256, key="width_preprocess")
            params["height"] = st.slider("Height", 50, 1000, 256, key="height_preprocess")
        elif operation == "rotate":
            params["angle"] = st.slider("Angle", 0, 360, 90, key="angle_preprocess")
        elif operation == "blur":
            params["radius"] = st.slider("Blur Radius", 1, 10, 2, key="radius_preprocess")
        elif operation == "threshold":
            params["threshold"] = st.slider("Threshold", 0, 255, 128, key="thresh_preprocess")

        cols = st.columns(min(3, len(images)))
        for idx, img in enumerate(images):
            analyzer = ImageAnalyzer(img)
            processed_img = analyzer.preprocess(operation, **params)
            with cols[idx % len(cols)]:
                st.image(processed_img, caption=f"Image {idx + 1} ({operation})", use_container_width=True)
        tracker.log_operation(f"Applied {operation} preprocessing to selected images")

    with tab2:
        tracker.log_tab("Statistics")
        st.subheader("Image Statistics")
        for idx, img in enumerate(images):
            analyzer = ImageAnalyzer(img)
            stats = analyzer.get_stats()
            
            st.subheader(f"Image {idx + 1}")
            col1, col2 = st.columns(2)
            with col1:
                st.json(stats)
            with col2:
                fig = go.Figure()
                if isinstance(stats["mean"], list):
                    fig.add_trace(go.Bar(
                        x=['Red', 'Green', 'Blue'],
                        y=stats["mean"],
                        name='Mean',
                        marker_color=['red', 'green', 'blue']
                    ))
                else:
                    fig.add_trace(go.Bar(
                        x=['Gray'],
                        y=[stats["mean"]],
                        name='Mean',
                        marker_color=['gray']
                    ))
                st.plotly_chart(fig, use_container_width=True, key=f"stats_chart_{idx}")
        tracker.log_operation("Displayed image statistics for selected images")
    with tab3:
        tracker.log_tab("OCR")
        st.subheader("Text Extraction (OCR)")
        if st.button("Extract Text from Selected Images"):
            for idx, img in enumerate(images):
                analyzer = ImageAnalyzer(img)
                text = analyzer.extract_text()
                
                st.subheader(f"Image {idx + 1}")
                if text.strip():
                    st.text_area(f"Extracted Text from Image {idx + 1}", value=text, height=200)
                else:
                    st.warning(f"No text found in Image {idx + 1}")
        tracker.log_operation("Performed OCR on selected images")

    with tab4:
        tracker.log_tab("Visualizations")
        st.subheader("Advanced Visualizations")
        for idx, img in enumerate(images):
            analyzer = ImageAnalyzer(img)
            
            st.subheader(f"Image {idx + 1}")
            # Convert to grayscale if not already
            gray_img = analyzer.preprocess("grayscale")
            np_gray = np.array(gray_img)

            # Edge detection - using a colored colormap
            edges = feature.canny(np_gray, sigma=2)
            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 5))
            ax1.imshow(np_gray, cmap='viridis')  # Changed to colored colormap
            ax1.set_title('Grayscale (Colored)')
            ax2.imshow(edges, cmap='plasma')  # Changed to colored colormap
            ax2.set_title('Edge Detection (Colored)')
            st.pyplot(fig)

            # Histogram equalization - adding color to the histograms
            equalized = exposure.equalize_hist(np_gray)
            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 5))
            ax1.hist(np_gray.ravel(), bins=256, color='skyblue', edgecolor='navy')
            ax1.set_title('Original Histogram')
            ax2.hist(equalized.ravel(), bins=256, color='salmon', edgecolor='darkred')
            ax2.set_title('Equalized Histogram')
            st.pyplot(fig)
        tracker.log_operation("Viewed advanced visualizations")
        

    with tab5:
        tracker.log_tab("Metadata")
        st.subheader("Image Metadata")
        cols = st.columns(min(3, len(images)))
        for idx, img in enumerate(images):
            analyzer = ImageAnalyzer(img)
            with cols[idx % len(cols)]:
                st.subheader(f"Image {idx + 1}")
                st.json(analyzer.metadata)
                st.image(analyzer.get_histogram(), caption="Color Histogram", use_container_width=True)  
        tracker.log_operation("Viewed image metadata")

def extraction_ui():
    tracker = st.session_state.tracker  # Ensure tracker is available
    session_manager = get_session_manager()  # Get session manager instance
    section = "Extraction"  # Define section name

    tracker.log_section(section)
    """Streamlit UI for the extraction feature with advanced image analysis"""
    st.header("Universal Data Extraction with Advanced Image Analysis")
    
    # Use session manager for file upload
    uploaded_filext = handle_file_upload(
        section=section,
        file_types=["pdf", "docx", "pptx", "xlsx", "xls", "csv"],
        title="Upload documents (PDF, DOCX, PPTX, Excel, CSV)",
        help_text="Supported formats: PDF, Word, PowerPoint, Excel, CSV"
    )
    
    # Display file info if processed
    if session_manager.get_data(section, 'file_processed', False):
        display_file_info(section)
    
    # Check if we have an uploaded file to process
    if uploaded_filext is not None:
        with st.spinner("Extracting content..."):
            text, images = universal_extraction(uploaded_filext)
            tracker.log_operation(f"Extracted content from {uploaded_filext.name}")
            
            if text or images:
                # Store results in session manager
                session_manager.set_data(section, 'extraction_result', {
                    "text": text,
                    "images": images,
                    "file_name": uploaded_filext.name
                })
    
    # Get extraction results from session manager
    extraction_result = session_manager.get_data(section, 'extraction_result')
    
    if extraction_result:
        st.subheader(f"Extracted from: {extraction_result['file_name']}")
        
        # Main tabs for different sections
        tab_text, tab_images, tab_analysis = st.tabs([
            "Text Content", "Extracted Images", "Image Analysis"
        ])
        
        with tab_text:
            tracker.log_tab("Text Content")
            if extraction_result['text']:
                st.subheader("Extracted Text")
                st.text_area("Text Content", 
                           value=extraction_result['text'], 
                           height=500,
                           key=f"{section}_text_area")
                tracker.log_operation("Displayed extracted text")
            else:
                st.warning("No text found in document")
                
        with tab_images:
            tracker.log_tab("Extracted Images")
            if extraction_result['images']:
                st.subheader(f"All Extracted Images ({len(extraction_result['images'])})")
                
                # Show all images in a grid
                cols = st.columns(3)
                for idx, img in enumerate(extraction_result['images']):
                    with cols[idx % 3]:
                        st.image(img, 
                               caption=f"Image {idx + 1}", 
                               use_container_width=True)
                        
                        # Download button for each image
                        buf = io.BytesIO()
                        img.save(buf, format='PNG')
                        buf.seek(0)
                        st.download_button(
                            label=f"Download Image {idx + 1}",
                            data=buf,
                            file_name=f"image_{idx + 1}.png",
                            mime="image/png",
                            key=f"{section}_download_{idx}"
                        )
                tracker.log_operation("Displayed and enabled download for extracted images")
                                
                # Image clustering analysis
                if len(extraction_result['images']) > 1:
                    st.subheader("Image Similarity Analysis")
                    if st.button("Cluster Similar Images", key=f"{section}_cluster_btn"):
                        reduced, clusters = cluster_images(extraction_result['images'])
                        if len(reduced) > 0:
                            fig = px.scatter(
                                x=reduced[:, 0], 
                                y=reduced[:, 1],
                                color=clusters,
                                title="Image Clustering (2D PCA Projection)"
                            )
                            st.plotly_chart(fig, use_container_width=True)
                            
                            # Show cluster distribution
                            cluster_counts = Counter(clusters)
                            fig = px.bar(
                                x=[f"Cluster {k}" for k in cluster_counts.keys()],
                                y=list(cluster_counts.values()),
                                title="Cluster Distribution"
                            )
                            st.plotly_chart(fig, use_container_width=True)
                            tracker.log_operation("Ran image clustering analysis")
            else:
                st.info("No images found in document")
        
        with tab_analysis:
            tracker.log_tab("Image Analysis")
            if extraction_result['images']:
                st.subheader("Image Analysis")
                selected_indices = st.multiselect(
                    "Select images to analyze",
                    options=[f"Image {i+1}" for i in range(len(extraction_result['images']))],
                    default=[f"Image {i+1}" for i in range(min(3, len(extraction_result['images'])))],
                    key=f"{section}_image_select"
                )
                
                if selected_indices:
                    selected_images = [extraction_result['images'][int(idx.split()[1])-1] for idx in selected_indices]
                    tracker.log_operation(f"Selected images for analysis: {selected_indices}")
                    show_image_analysis(selected_images)
            else:
                st.info("No images available for analysis")
    
    if st.button("Clear Extraction Results", key=f"{section}_clear_btn"):
        session_manager.clear_section_data(section)
        st.rerun()