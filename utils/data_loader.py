"""
Enhanced data loader with section isolation
Handles loading of various file formats with optimized memory usage
and prevents data leakage between sections
"""

import streamlit as st
import pandas as pd
import numpy as np
import tempfile
import os
import fitz  # PyMuPDF
import io
from PIL import Image
import chardet
import re
import psutil
import time
import base64
import concurrent.futures
from typing import Dict, Any, Optional, Tuple, List, Union
import dask.dataframe as dd
from io import BytesIO
from langdetect import detect
from langdetect.lang_detect_exception import LangDetectException
import PyPDF2
from utils.session_state_manager import get_session_manager

# Get the session manager instance
session_manager = get_session_manager()

def detect_file_encoding(file):
    """
    Detect the encoding of a file
    Args:
        file: File object
    Returns:
        Detected encoding
    """
    raw_data = file.read()
    result = chardet.detect(raw_data)
    file.seek(0)  # Reset file pointer to the beginning
    return result['encoding']

@st.cache_data(show_spinner=False)
def load_data(file, section: str):
    """
    Load data from a file into section-specific session state
    Args:
        file: Uploaded file object
        section: Section name where file is being loaded
    Returns:
        Loaded data (DataFrame, text, or image)
    """
    try:
        # Create a temporary file for all formats
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_path = tmp_file.name
            
        if file.name.endswith('.csv'):
            # Efficient CSV loading with optimized parameters
            df = pd.read_csv(tmp_path, engine='c', low_memory=False)
            os.unlink(tmp_path)
            return df
        elif file.name.endswith('.xlsx') or file.name.endswith('.xls'):
            # Efficient Excel loading
            df = pd.read_excel(tmp_path, engine='openpyxl')
            os.unlink(tmp_path)
            return df
        elif file.name.endswith('.pdf'):
            # Extract text from PDF and return as a string (not a DataFrame)
            with open(tmp_path, 'rb') as f:
                text, images = extract_text_and_images_from_pdf(f)
            os.unlink(tmp_path)
            return {'text': text, 'images': images}
        elif file.name.endswith('.txt'):
            # Read plain text file efficiently
            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            os.unlink(tmp_path)
            return text
        elif file.name.endswith('.docx'):
            # Process DOCX files
            try:
                from docx import Document
                doc = Document(tmp_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                os.unlink(tmp_path)
                return text
            except ImportError:
                st.error("To process DOCX files, please install python-docx: pip install python-docx")
                os.unlink(tmp_path)
                return None
        elif file.name.endswith('.pptx'):
            # Process PPTX files
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                text_data = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_data.append(shape.text)
                text = "\n".join(text_data)
                os.unlink(tmp_path)
                return text
            except ImportError:
                st.error("To process PPTX files, please install python-pptx: pip install python-pptx")
                os.unlink(tmp_path)
                return None
        elif file.name.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif', '.webp')):
            # Process image files
            image = Image.open(tmp_path)
            os.unlink(tmp_path)
            return image
        else:
            st.error(f"Unsupported file format in {section}. Please upload a supported file type.")
            os.unlink(tmp_path)
            return None
    except Exception as e:
        st.error(f"Error loading the dataset in {section}: {e}")
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def load_large_data(file, section: str, file_type=None, blocksize="64MB", sample_rows=10000):
    """
    Efficiently loads large datasets using optimized techniques for a specific section
    
    Args:
        file: Uploaded file object
        section: Section name where file is being loaded
        file_type: Type of file (csv, excel, txt, pdf, docx, pptx)
        blocksize: Block size for chunked reading
        sample_rows: Number of rows to sample for preview
        
    Returns:
        Tuple of (complete_data, preview_data)
    """
    try:
        # Create temporary file for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_path = tmp_file.name
            
        # Get file size for optimization decisions
        file_size = os.path.getsize(tmp_path) / (1024 * 1024)  # Size in MB
        st.info(f"File size: {file_size:.2f} MB - Optimizing loading process...")
        
        # Start resource monitoring
        start_time = time.time()
        mem_before = psutil.virtual_memory().used / (1024 * 1024)
        
        # Infer file type if not provided
        if file_type is None:
            fname = file.name.lower()
            if fname.endswith('.csv'):
                file_type = 'csv'
            elif fname.endswith('.xlsx') or fname.endswith('.xls'):
                file_type = 'excel'
            elif fname.endswith('.txt'):
                file_type = 'txt'
            elif fname.endswith('.pdf'):
                file_type = 'pdf'
            elif fname.endswith('.docx'):
                file_type = 'docx'
            elif fname.endswith('.pptx'):
                file_type = 'pptx'
            else:
                st.error(f"Unsupported file format for large data loading in {section}.")
                os.unlink(tmp_path)
                return None, None

        # Create progress bar
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        def update_progress(step, total_steps, message):
            progress = int((step / total_steps) * 100)
            progress_bar.progress(progress)
            status_text.text(f"{message} ({progress}%)")
        
        # CSV - Optimized loading with Dask and parallel processing
        if file_type == 'csv':
            update_progress(1, 4, "Analyzing CSV structure")
            
            # Optimize chunk size based on file size
            optimized_blocksize = "256MB" if file_size > 500 else "64MB"
            
            update_progress(2, 4, "Loading data in parallel chunks")
            df = dd.read_csv(
                tmp_path, 
                blocksize=optimized_blocksize,
                assume_missing=True,
                dtype=str,
                sample=1000000  # Optimized sampling
            )
            
            update_progress(3, 4, "Preparing preview")
            preview = df.head(sample_rows, compute=True)
            
            update_progress(4, 4, "Finalizing dataset")
            st.info(f"Previewing first {sample_rows} rows. Full dataset will be processed in chunks.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return df, preview
        
        # Excel - Optimized loading with chunk processing
        elif file_type == 'excel':
            update_progress(1, 3, "Reading Excel file")
            
            # For large files, read in chunks
            if file_size > 50:  # Files larger than 50MB
                chunks = []
                chunk_size = 10000
                total_rows = 0
                
                # Read Excel in chunks
                reader = pd.read_excel(tmp_path, engine='openpyxl', chunksize=chunk_size)
                for i, chunk_df in enumerate(reader):
                    chunks.append(chunk_df)
                    total_rows += len(chunk_df)
                    update_progress(i+1, int(2000000/chunk_size), f"Processing chunk {i+1}")
                
                df = pd.concat(chunks, ignore_index=True)
            else:
                df = pd.read_excel(tmp_path, engine='openpyxl')
            
            update_progress(2, 3, "Optimizing data types")
            df = optimize_data_types(df)
            
            update_progress(3, 3, "Finalizing dataset")
            st.info(f"Loaded {len(df)} rows from Excel file.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return df, df.head(sample_rows)
        
        # Text - Efficient reading with progress
        elif file_type == 'txt':
            update_progress(1, 2, "Reading text file")
            
            # Read large text files in chunks
            text = ""
            chunk_size = 1024 * 1024  # 1MB chunks
            total_chunks = (os.path.getsize(tmp_path) // chunk_size) + 1
            
            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                for i in range(total_chunks):
                    chunk = f.read(chunk_size)
                    text += chunk
                    update_progress(i+1, total_chunks, "Reading text chunks")
            
            update_progress(2, 2, "Finalizing text")
            st.info("Loaded plain text file.")
            
            # Resource usage report
            mem_after = psutil.virtual_memory().used / (1024 * 1024)
            load_time = time.time() - start_time
            st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
            
            os.unlink(tmp_path)
            return text, text[:10000]  # Return full text and preview
        
        # PDF - Parallel processing of pages
        elif file_type == 'pdf':
            update_progress(1, 3, "Extracting PDF content")
            with open(tmp_path, 'rb') as f:
                text, images = extract_text_and_images_from_pdf(f)
            
            update_progress(2, 3, "Processing extracted content")
            if text or images:
                st.info("Loaded PDF file.")
                
                # Resource usage report
                mem_after = psutil.virtual_memory().used / (1024 * 1024)
                load_time = time.time() - start_time
                st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
                
                os.unlink(tmp_path)
                return {'text': text, 'images': images}, {'text': text[:10000], 'images': images[:5] if images else []}
            else:
                st.error("No content could be extracted from the PDF file.")
                os.unlink(tmp_path)
                return None, None
        
        # DOCX - Process Word documents
        elif file_type == 'docx':
            update_progress(1, 2, "Processing DOCX file")
            try:
                from docx import Document
                doc = Document(tmp_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                
                update_progress(2, 2, "Finalizing document text")
                st.info("Loaded DOCX file.")
                
                # Resource usage report
                mem_after = psutil.virtual_memory().used / (1024 * 1024)
                load_time = time.time() - start_time
                st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
                
                os.unlink(tmp_path)
                return text, text[:10000]  # Return full text and preview
            except ImportError:
                st.error("To process DOCX files, please install python-docx: pip install python-docx")
                os.unlink(tmp_path)
                return None, None
        
        # PPTX - Process PowerPoint documents
        elif file_type == 'pptx':
            update_progress(1, 2, "Processing PPTX file")
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                text_data = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_data.append(shape.text)
                text = "\n".join(text_data)
                
                update_progress(2, 2, "Finalizing presentation text")
                st.info("Loaded PPTX file.")
                
                # Resource usage report
                mem_after = psutil.virtual_memory().used / (1024 * 1024)
                load_time = time.time() - start_time
                st.success(f"Loaded in {load_time:.2f} seconds | Memory used: {mem_after - mem_before:.2f} MB")
                
                os.unlink(tmp_path)
                return text, text[:10000]  # Return full text and preview
            except ImportError:
                st.error("To process PPTX files, please install python-pptx: pip install python-pptx")
                os.unlink(tmp_path)
                return None, None
        
        else:
            st.error(f"Only CSV, Excel, PDF, TXT, DOCX, PPTX files are supported for large data loading in {section}.")
            os.unlink(tmp_path)
            return None, None
    
    except Exception as e:
        st.error(f"Error loading large dataset in {section}: {e}")
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None, None
    finally:
        if 'progress_bar' in locals():
            progress_bar.empty()
        if 'status_text' in locals():
            status_text.empty()

def is_structured(df):
    """
    Enhanced structured data check with additional validations
    
    Args:
        df: DataFrame to check
        
    Returns:
        Boolean indicating if data is structured
    """
    try:
        # Basic emptiness check
        if df.empty or df.isnull().all().all():
            return False
        
        # Check for minimum structure requirements
        if len(df.columns) < 1 or df.shape[0] < 1:
            return False

        # Check for consistent number of columns across all rows
        col_count = len(df.columns)
        if any(len(row) != col_count for row in df.itertuples(index=False)):
            return False

        # Check for valid header (first row should contain mostly strings)
        header_validity = sum(isinstance(col, str) for col in df.columns) / len(df.columns)
        if header_validity < 0.8:  # At least 80% of headers should be strings
            return False

        # Check for non-empty data cells
        non_empty_cells = df.count().sum()
        if non_empty_cells / (df.shape[0] * df.shape[1]) < 0.1:  # At least 10% data present
            return False

        return True
    except Exception as e:
        st.error(f"Error checking dataset structure: {e}")
        return False

def optimize_data_types(df):
    """
    Optimize DataFrame memory usage by converting to appropriate data types
    
    Args:
        df: DataFrame to optimize
        
    Returns:
        Optimized DataFrame
    """
    for col in df.columns:
        if df[col].dtype == 'float64':
            df[col] = df[col].astype('float32')
        elif df[col].dtype == 'int64':
            df[col] = df[col].astype('int32')
        elif df[col].dtype == 'object':
            if df[col].nunique() / len(df[col]) < 0.5:  # Convert to category if unique values are less than 50%
                df[col] = df[col].astype('category')
    return df

def transform_unstructured_data(df):
    """
    Transform unstructured data into structured data
    
    Args:
        df: DataFrame to transform
        
    Returns:
        Transformed DataFrame
    """
    try:
        # Step 1: Drop completely empty rows and columns
        df_cleaned = df.dropna(how='all').dropna(axis=1, how='all')
        
        # Step 2: Fill missing values with placeholders
        df_cleaned = df_cleaned.fillna('MISSING_VALUE')
        
        # Step 3: Identify header row (row with most non-empty values)
        non_empty_counts = df_cleaned.apply(lambda row: row.astype(bool).sum(), axis=1)
        if non_empty_counts.max() > 1:
            header_row_idx = non_empty_counts.idxmax()
            new_header = df_cleaned.iloc[header_row_idx]
            df_cleaned = df_cleaned.drop(header_row_idx)
            df_cleaned.columns = new_header
        
        # Step 4: Reset index
        df_cleaned = df_cleaned.reset_index(drop=True)
        
        # Step 5: Clean column names
        df_cleaned.columns = [re.sub(r'\W+', '_', str(col)).strip('_') for col in df_cleaned.columns]
        
        return df_cleaned
    except Exception as e:
        st.error(f"Error transforming unstructured data: {e}")
        return df

def extract_text_and_images_from_pdf(uploaded_pdf):
    """
    Extract text and images from PDF
    
    Args:
        uploaded_pdf: PDF file object
        
    Returns:
        Tuple of (text, images)
    """
    text = ""
    images = []
    try:
        # First try text extraction
        try:
            reader = PyPDF2.PdfReader(uploaded_pdf)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            st.warning(f"Standard text extraction failed: {e}")
        
        # If no text found, try OCR fallback
        if not text.strip():
            st.info("No text found. Attempting image extraction...")
            uploaded_pdf.seek(0)
            images = extract_images_from_pdf(uploaded_pdf)
            
            # If no images either, try OCR on the PDF pages
            if not images:
                st.info("No images found. Trying alternative text extraction...")
                uploaded_pdf.seek(0)
                with fitz.open(stream=uploaded_pdf.read(), filetype="pdf") as doc:
                    for page in doc:
                        text += page.get_text() or ""
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
    
    return text, images

def extract_images_from_pdf(pdf_file):
    """
    Extract images from PDF
    
    Args:
        pdf_file: PDF file object
        
    Returns:
        List of images
    """
    images = []
    pdf_file.seek(0)
    try:
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page_index in range(len(doc)):
            for img_index, img in enumerate(doc.get_page_images(page_index)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)
    except Exception as e:
        st.error(f"Error extracting images from PDF: {e}")
    return images

def handle_file_upload(section: str, file_types=None, title="Upload a file", help_text=None):
    """
    Handle file upload for a specific section with proper isolation
    
    Args:
        section: Section name
        file_types: List of accepted file extensions
        title: Title for the file uploader
        help_text: Help text for the file uploader
        
    Returns:
        Uploaded file or None
    """
    if file_types is None:
        file_types = ["csv", "xlsx", "xls", "txt", "pdf", "png", "jpg", "jpeg", "bmp", "tiff", "gif", "webp", "docx", "pptx"]
    
    # Get a unique file uploader key for this section
    uploader_key = session_manager.get_file_uploader_key(section)
    
    # Display the file uploader
    uploaded_file = st.file_uploader(
        title,
        type=file_types,
        key=f"{section}_{uploader_key}",
        help=help_text
    )
    
    # Process the uploaded file
    if uploaded_file is not None:
        # Store file in section-specific session state
        session_manager.set_data(section, 'uploaded_file', uploaded_file)
        session_manager.set_data(section, 'uploaded_file_name', uploaded_file.name)
        session_manager.set_data(section, 'file_processed', False)
        
        # Process file based on type
        try:
            if uploaded_file.name.endswith(('.txt', '.docx', '.pptx')):
                # Text-based files
                if uploaded_file.name.endswith('.txt'):
                    text_data = uploaded_file.read().decode("utf-8", errors="ignore")
                elif uploaded_file.name.endswith('.docx'):
                    try:
                        from docx import Document
                        doc = Document(uploaded_file)
                        text_data = "\n".join([para.text for para in doc.paragraphs])
                    except ImportError:
                        st.error("To process DOCX files, please install python-docx: pip install python-docx")
                        text_data = ""
                elif uploaded_file.name.endswith('.pptx'):
                    try:
                        from pptx import Presentation
                        prs = Presentation(uploaded_file)
                        text_data = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_data.append(shape.text)
                        text_data = "\n".join(text_data)
                    except ImportError:
                        st.error("To process PPTX files, please install python-pptx: pip install python-pptx")
                        text_data = ""
                
                session_manager.set_data(section, 'text_data', text_data)
                session_manager.set_data(section, 'file_processed', True)
                st.success(f"‚úÖ {uploaded_file.name} uploaded successfully to {section}!")
                
            elif uploaded_file.name.endswith('.pdf'):
                # PDF files
                text, images = extract_text_and_images_from_pdf(uploaded_file)
                session_manager.set_data(section, 'pdf_data', {
                    'text': text,
                    'images': images
                })
                session_manager.set_data(section, 'file_processed', True)
                
                if text or images:
                    success_msg = f"PDF file uploaded successfully to {section}!"
                    if text:
                        success_msg += " Text extracted."
                    if images:
                        success_msg += f" {len(images)} images extracted."
                    st.success(f"‚úÖ {success_msg}")
                else:
                    st.warning(f"PDF uploaded to {section} but no content could be extracted.")
            
            elif uploaded_file.type.startswith("image/"):
                # Image files
                image = Image.open(uploaded_file)
                session_manager.set_data(section, 'image_data', image)
                session_manager.set_data(section, 'file_processed', True)
                st.success(f"‚úÖ Image file uploaded successfully to {section}!")
            
            else:
                # Tabular data (CSV, Excel)
                df = load_data(uploaded_file, section)
                if df is not None:
                    # Check if structured
                    if not is_structured(df):
                        st.warning(f"Dataset appears unstructured in {section}. Transforming...")
                        df = transform_unstructured_data(df)
                        
                        # Check again after transformation
                        if is_structured(df):
                            st.success(f"Dataset successfully transformed to structured format in {section}!")
                        else:
                            st.error(f"Could not transform to structured format in {section}. Proceeding with original data.")
                    
                    df = optimize_data_types(df)
                    session_manager.set_data(section, 'df', df)
                    session_manager.set_data(section, 'is_structured', is_structured(df))
                    session_manager.set_data(section, 'file_processed', True)
                    st.success(f"‚úÖ Dataset is ready for analysis in {section}!")
                else:
                    st.error(f"‚ùå Error: The uploaded file could not be loaded in {section}.")
        except Exception as e:
            st.error(f"‚ùå Error processing file in {section}: {str(e)}")
            session_manager.set_data(section, 'file_processed', False)
        
        return uploaded_file
    
    return None

def clear_section_data(section: str):
    """
    Clear all data for a specific section
    
    Args:
        section: Section name
    """
    session_manager.clear_section_data(section)
    st.success(f"All data has been cleared from {section}!")

def display_file_info(section: str):
    """
    Display file information for a specific section
    
    Args:
        section: Section name
    """
    file_info = session_manager.get_file_info(section)
    
    if file_info['file_processed']:
        st.markdown(f"**Uploaded File:** `{file_info['uploaded_file_name']}`")
        
        # Add a remove button
        if st.button("üóëÔ∏è Remove Uploaded File", key=f"remove_file_btn_{section}"):
            clear_section_data(section)
            st.rerun()

def plot_data_summary(section: str):
    """
    Comprehensive data summary with statistics and visual overview
    
    Args:
        section: Section name
    """
    try:
        df = session_manager.get_dataframe(section)
        if df is None or not hasattr(df, "columns") or df.empty:
            st.warning(f"Please upload and load a dataset first in {section}. No data available for summary.")
            return

        summary = []
        for col in df.columns:
            dtype = str(df[col].dtype)
            unique = df[col].nunique()
            missing = df[col].isnull().sum()
            # Convert sample values to string to avoid Arrow conversion issues
            sample = str(df[col].dropna().unique()[:5])
            summary.append({
                "Column": col,
                "Type": dtype,
                "Unique Values": unique,
                "Missing Values": missing,
                "Sample Values": sample
            })
        summary_df = pd.DataFrame(summary)
        st.dataframe(summary_df)
    except Exception as e:
        st.error(f"Error in data summary for {section}: {e}")
        return None